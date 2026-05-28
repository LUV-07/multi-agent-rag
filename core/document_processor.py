"""
core/document_processor.py — Universal document parser
Supports: PDF, DOCX, TXT, MD, PPTX, XLSX, HTML

"""
import os
from typing import List, Dict, Optional
from utils.logger import get_logger

log = get_logger("document_processor")


class DocumentProcessor:
    """
    Extracts raw text from various document formats.
    Returns list of page dicts: {text, page_number, source}
    """

    def process(self, file_path: str) -> List[Dict]:
        """Auto-detect format and extract text."""
        ext = os.path.splitext(file_path)[1].lower()
        source_name = os.path.basename(file_path)

        log.info(f"Processing {source_name} (type: {ext})")

        processors = {
            ".pdf":  self._process_pdf,
            ".docx": self._process_docx,
            ".txt":  self._process_txt,
            ".md":   self._process_txt,
            ".pptx": self._process_pptx,
            ".xlsx": self._process_xlsx,
            ".html": self._process_html,
            ".htm":  self._process_html,
        }

        processor = processors.get(ext)
        if not processor:
            raise ValueError(f"Unsupported file type: {ext}")

        pages = processor(file_path)

        # Attach source name to all pages
        for page in pages:
            page["source"] = source_name

        log.info(f"Extracted {len(pages)} pages/sections from {source_name}")
        return pages

    #------PDF Extraction------

    def _process_pdf(self, file_path: str) -> List[Dict]:
        """Extract text from PDF using PyMuPDF (fitz)."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("Install PyMuPDF: pip install pymupdf")

        pages = []
        doc = fitz.open(file_path)

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                pages.append({
                    "text": text,
                    "page": page_num,
                    "total_pages": len(doc),
                })

        doc.close()
        return pages

    #------DOCX Extraction------

    def _process_docx(self, file_path: str) -> List[Dict]:
        """Extract text from Word document."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("Install python-docx: pip install python-docx")

        doc = Document(file_path)
        full_text = []

        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)

        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    full_text.append(row_text)

        combined = "\n".join(full_text)
        return [{"text": combined, "page": 1, "total_pages": 1}]

    #------TXT / MD Extraction------

    def _process_txt(self, file_path: str) -> List[Dict]:
        """Extract text from plain text or markdown file."""
        encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]

        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    text = f.read()
                return [{"text": text, "page": 1, "total_pages": 1}]
            except UnicodeDecodeError:
                continue

        raise ValueError("Could not decode text file with any supported encoding.")

    #-----PPTX Extraction------

    def _process_pptx(self, file_path: str) -> List[Dict]:
        """Extract text from PowerPoint slides."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Install python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        pages = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                pages.append({
                    "text": "\n".join(slide_text),
                    "page": slide_num,
                    "total_pages": len(prs.slides),
                })

        return pages

    #------ XLSX Extarction ------

    def _process_xlsx(self, file_path: str) -> List[Dict]:
        """Extract text from Excel spreadsheet."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("Install openpyxl: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, data_only=True)
        pages = []

        for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            rows_text = []

            for row in ws.iter_rows(values_only=True):
                row_str = " | ".join(str(cell) for cell in row if cell is not None and str(cell).strip())
                if row_str:
                    rows_text.append(row_str)

            if rows_text:
                pages.append({
                    "text": f"Sheet: {sheet_name}\n" + "\n".join(rows_text),
                    "page": sheet_num,
                    "total_pages": len(wb.sheetnames),
                })

        return pages

    #------ HTML Extraction ------

    def _process_html(self, file_path: str) -> List[Dict]:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("Install beautifulsoup4: pip install beautifulsoup4")

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "lxml")

        # Remove scripts and styles
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        return [{"text": text, "page": 1, "total_pages": 1}]
